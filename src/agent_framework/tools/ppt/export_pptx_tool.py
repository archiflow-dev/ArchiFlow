"""
Export PPTX Tool for PPT Agent MVP.

This tool converts generated slide images into PowerPoint presentations.
It finds all slide_*.png files in the current directory and creates a PPTX
with proper 16:9 aspect ratio (10" x 5.625").
"""

import os
import logging
import glob
from datetime import datetime
from typing import Optional, Dict, Any, List, ClassVar
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    Presentation = None
    Inches = None
    Emu = None

from ..tool_base import BaseTool, ToolResult
from PIL import Image

# Set up logger for this module
logger = logging.getLogger(__name__)


class ExportPPTXTool(BaseTool):
    """
    Tool for exporting slide images to PowerPoint presentation format.

    This tool finds all slide_*.png files in the current directory
    and creates a PowerPoint presentation with proper formatting.
    """

    name: str = "export_pptx"
    description: str = "Export generated slide images to a PowerPoint presentation file."

    parameters: Dict[str, Any] = {
        "type": "object",
        "properties": {
            "title": {
                "type": "string",
                "description": "Title for the presentation"
            },
            "input_dir": {
                "type": "string",
                "description": "Input directory containing slide images (default: data/sessions/{session_id}/images)"
            },
            "output_dir": {
                "type": "string",
                "description": "Output directory for the PPTX file (default: data/sessions/{session_id}/ppt_exports)"
            },
            "slide_pattern": {
                "type": "string",
                "description": "Pattern to match slide files (default: slide_*.png)",
                "default": "slide_*.png"
            }
        },
        "required": ["title"]
    }

    # PowerPoint dimensions for 16:9 aspect ratio
    SLIDE_WIDTH_INCHES: ClassVar[float] = 10.0
    SLIDE_HEIGHT_INCHES: ClassVar[float] = 5.625

    def __init__(self, **data):
        """Initialize the ExportPPTXTool."""
        super().__init__(**data)

        if not PPTX_AVAILABLE:
            logger.warning(
                "python-pptx package is not installed. "
                "Install with: pip install python-pptx"
            )

    async def execute(
        self,
        title: str,
        input_dir: Optional[str] = None,
        session_id: str = "default",
        output_dir: Optional[str] = None,
        slide_pattern: str = "slide_*.png",
        **kwargs
    ) -> ToolResult:
        """
        Export slide images to a PowerPoint presentation.

        Args:
            title: Title for the presentation
            input_dir: Input directory containing slide images (default: data/sessions/{session_id}/images)
            session_id: Session ID for organizing files
            output_dir: Output directory for the PPTX file (default: data/sessions/{session_id}/ppt_exports)
            slide_pattern: Pattern to match slide files

        Returns:
            ToolResult containing the file path and slide count or error message
        """
        if not PPTX_AVAILABLE:
            return self.fail_response(
                "python-pptx package is required. Install with: pip install python-pptx"
            )

        try:
            logger.info(f"Starting PPTX export for presentation: {title}")

            # Determine input directory
            if input_dir is None:
                # Default to data/sessions/{session_id}/images
                input_dir = os.path.join("data", "sessions", session_id, "images")

            # Determine output directory
            if output_dir is None:
                # Default to data/sessions/{session_id}/ppt_exports
                output_dir = os.path.join("data", "sessions", session_id, "ppt_exports")

            # Validate input directory
            input_path = Path(input_dir)
            if not input_path.exists():
                return self.fail_response(
                    f"Input directory does not exist: {input_dir}"
                )

            if not input_path.is_dir():
                return self.fail_response(
                    f"Input path is not a directory: {input_dir}"
                )

            # Find all slide files
            slide_files = self._find_slide_files(slide_pattern, input_path)
            if not slide_files:
                return self.fail_response(
                    f"No slide files found matching pattern '{slide_pattern}' in directory: {input_dir}"
                )

            logger.info(f"Found {len(slide_files)} slide files")

            # Create output directory if it doesn't exist
            output_path = Path(output_dir)
            output_path.mkdir(parents=True, exist_ok=True)
            logger.info(f"Output directory: {output_path.absolute()}")

            # Generate filename with timestamp
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            # Sanitize title for filename
            safe_title = "".join(c for c in title if c.isalnum() or c in (' ', '-', '_')).rstrip()
            safe_title = safe_title.replace(' ', '_')
            filename = f"{safe_title}_{timestamp}.pptx"
            filepath = output_path / filename

            # Create PowerPoint presentation
            prs = Presentation()

            # Set slide size to 16:9
            prs.slide_width = Inches(self.SLIDE_WIDTH_INCHES)
            prs.slide_height = Inches(self.SLIDE_HEIGHT_INCHES)

            # Add slides
            slide_count = 0
            for slide_file in slide_files:
                try:
                    # Load the image to get dimensions
                    img_path = Path(slide_file)
                    if not img_path.exists():
                        logger.warning(f"Slide file not found: {slide_file}")
                        continue

                    with Image.open(img_path) as img:
                        img_width, img_height = img.size

                    # Add a blank slide
                    slide_layout = prs.slide_layouts[6]  # Blank layout
                    slide = prs.slides.add_slide(slide_layout)

                    # Calculate image dimensions to fit slide
                    slide_width_emu = prs.slide_width
                    slide_height_emu = prs.slide_height

                    # Determine scaling to fit
                    width_scale = slide_width_emu / (img_width * 9525)  # Convert pixels to EMU
                    height_scale = slide_height_emu / (img_height * 9525)
                    scale = min(width_scale, height_scale)

                    # Calculate final dimensions
                    final_width = img_width * 9525 * scale
                    final_height = img_height * 9525 * scale

                    # Center the image on the slide
                    left = (slide_width_emu - final_width) / 2
                    top = (slide_height_emu - final_height) / 2

                    # Add image to slide
                    slide.shapes.add_picture(
                        str(img_path.absolute()),
                        left=left,
                        top=top,
                        width=final_width,
                        height=final_height
                    )

                    slide_count += 1
                    logger.info(f"Added slide {slide_count}: {img_path.name}")

                except Exception as e:
                    logger.error(f"Error adding slide {slide_file}: {e}")
                    continue

            if slide_count == 0:
                return self.fail_response("No slides were successfully added to the presentation")

            # Save the presentation
            prs.save(str(filepath))
            logger.info(f"Presentation saved to: {filepath}")

            # Return success with file information
            result = {
                "success": True,
                "file_path": str(filepath.absolute()),
                "filename": filename,
                "title": title,
                "session_id": session_id,
                "slide_count": slide_count,
                "input_dir": str(input_path.absolute()),
                "output_dir": str(output_path.absolute()),
                "aspect_ratio": "16:9",
                "dimensions": f"{self.SLIDE_WIDTH_INCHES}\" x {self.SLIDE_HEIGHT_INCHES}\""
            }

            return self.success_response(result)

        except Exception as e:
            error_msg = f"Error exporting to PPTX: {str(e)}"
            logger.error(error_msg, exc_info=True)
            return self.fail_response(error_msg)

    def _find_slide_files(self, pattern: str, input_dir: Path = None) -> List[str]:
        """
        Find slide files matching the pattern.

        Args:
            pattern: Glob pattern for slide files
            input_dir: Directory to search in (defaults to current directory)

        Returns:
            Sorted list of slide file paths
        """
        if input_dir is None:
            input_dir = Path(".")

        # Use glob to find matching files in the input directory
        search_pattern = str(input_dir / pattern)
        files = glob.glob(search_pattern)

        # Sort numerically (slide_001.png, slide_002.png, etc.)
        def extract_number(f):
            import re
            match = re.search(r'slide_(\d+)', f)
            return int(match.group(1)) if match else 0

        files.sort(key=extract_number)
        return files

    def get_slide_count(self, pattern: str = "slide_*.png", input_dir: str = None, session_id: str = "default") -> int:
        """
        Get the count of slide files without creating a presentation.

        Args:
            pattern: Pattern to match slide files
            input_dir: Directory to search in (if None, uses data/sessions/{session_id}/images)
            session_id: Session ID for organizing files

        Returns:
            Number of slide files found
        """
        if input_dir is None:
            input_dir = os.path.join("data", "sessions", session_id, "images")
        return len(self._find_slide_files(pattern, Path(input_dir)))

    def __repr__(self):
        """String representation of the tool."""
        return f"ExportPPTXTool(pptx_available={PPTX_AVAILABLE})"